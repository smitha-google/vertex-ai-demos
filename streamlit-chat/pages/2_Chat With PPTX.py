import streamlit as st
from pptx import Presentation
from pptx.util import Inches
from typing import Dict, List, Union, Any, Iterable
import pandas as pd
from vertexai.preview.generative_models import (Content,
                                                GenerationConfig,
                                                GenerativeModel,
                                                GenerationResponse,
                                                Image, 
                                                HarmCategory, 
                                                HarmBlockThreshold, 
                                                Part)
from vertexai.language_models import TextEmbeddingModel
import vertexai
from google.cloud import aiplatform
from google.protobuf import struct_pb2
from smutils import utils

PROJECT_ID = "smithaargolisinternal"#os.environ.get('GCP_PROJECT') #Your Google Cloud Project ID
LOCATION = "us-central1"#os.environ.get('GCP_REGION')   #Your Google Cloud Project Region
vertexai.init(project=PROJECT_ID, location=LOCATION)
st.title("Chat With A Powerpoint")

if "messages" not in st.session_state:
    st.session_state.messages = []

#creating session states
if "text_df" not in st.session_state:
    st.session_state['text_df'] = pd.DataFrame()

if "image_df" not in st.session_state:
    st.session_state['image_df'] = pd.DataFrame()

if "image_uri" not in st.session_state:
    st.session_state['image_uri'] = ""

# Display chat messages from history on app rerun
for message in st.session_state.messages:
    with st.chat_message(message["role"]):
        st.markdown(message["content"])

if 'disabled' not in st.session_state:
    st.session_state.disabled = False

def disable():
    st.session_state.disabled = True

def load_ppt_as_dataframe(file:str):
    prs = Presentation("./pages/Doc AI Prep.pptx")
    text_df = pd.DataFrame(columns=['file_name', 'page_num', 'text', 'embeddings'])
    image_df = pd.DataFrame(columns=['file_name', 'page_num', 'text', 'embeddings'])
    image_uri= ""

    for slide_number, slide in enumerate(prs.slides): 
        st.write(slide_number)
        for shape in slide.shapes: 
            if hasattr(shape, "image"): 
                image_uri = utils.get_image_for_gemini_ppts(shape.image.blob, slide_number+1, "images", "Doc AI Prep",slide_number+1)
                image_embeddings = utils.get_image_embedding_from_multimodal_embedding_model(image_uri)
                df2 = {'file_name':"Doc AI Prep.pptx",'page_num': slide_number+1, 'text': shape.image, 'embeddings': image_embeddings} 
                image_df.loc[len(image_df)] = df2
                st.write(image_df)
                frames = [st.session_state['image_df'], image_df]
                st.session_state['image_df'] = pd.concat(frames)
                st.session_state['image_uri'] = image_uri
            elif hasattr(shape, "text"):
                text=shape.text[:1000]
                embeddings = utils.get_text_embedding_from_text_embedding_model(PROJECT_ID, text, 128)
                df2 = {'file_name':"Doc AI Prep.pptx", 'page_num': slide_number+1, 'text': text, 'embeddings': embeddings} 
                text_df.loc[len(text_df)] = df2
                frames = [st.session_state['text_df'], text_df]
                st.session_state['text_df'] = pd.concat(frames)
                st.write(text_df)
    return text_df

def get_metadata(file_name): 
   #Load the PPT file
   df = load_ppt_as_dataframe(file_name)

with st.form("my_form"):
    uploaded_ppt = st.file_uploader('Choose your .ppt file', type="pptx", 
                                disabled=st.session_state.disabled)
    submitted = st.form_submit_button("Load", on_click=disable())
    if submitted:
       st.info("successfully uploaded!")
       temp_file = "./temp.pptx"
       with open(temp_file, "wb") as file:
        file.write(uploaded_ppt.getvalue())
        file_name = uploaded_ppt.name
       st.info("Indexing and Generating Embeddings")
       get_metadata(file_name)
       st.info("File is ready for QnA")

@st.cache_resource
def load_model():
    model = GenerativeModel("gemini-pro")
    text_embedding_model = TextEmbeddingModel.from_pretrained("textembedding-gecko@001")
    return model

#Load the Gemini Model
model = load_model()

#Set the Config
config = {
        "temperature": 0.8,
        "max_output_tokens": 2048,
        }

def multiturn_generate_content(model: GenerativeModel,
                                  prompt: str, 
                                  generation_config: GenerationConfig,
                                  stream=True):
    safety_settings={
        HarmCategory.HARM_CATEGORY_HARASSMENT: HarmBlockThreshold.BLOCK_NONE,
        HarmCategory.HARM_CATEGORY_HATE_SPEECH: HarmBlockThreshold.BLOCK_NONE,
        HarmCategory.HARM_CATEGORY_SEXUALLY_EXPLICIT: HarmBlockThreshold.BLOCK_NONE,
        HarmCategory.HARM_CATEGORY_DANGEROUS_CONTENT: HarmBlockThreshold.BLOCK_NONE,
    }
    chat = model.start_chat()
    responses = chat.send_message(prompt, generation_config = generation_config,
                                      safety_settings=safety_settings,
                                      stream=True)

    final_response = []
    for response in responses:
        try:
            #st.write(response.text)
            final_response.append(response.text)
        except IndexError:
            # st.write(response)
            final_response.append("")
            continue
    return " ".join(final_response)

if prompt := st.chat_input("Add your prompt:"):
    st.chat_message("user").markdown(prompt)
    st.session_state.messages.append({"role": "user", "content": prompt})
    response=f"Echo: {prompt}"
    with st.spinner("Generating..."):
        matching_results_chunks_data = utils.get_similar_text_from_query(
            PROJECT_ID,prompt,st.session_state['text_df'],column_name="embeddings",
            top_n=3, chunk_text=False,
        )

        st.write(st.session_state['image_uri'])

        # Get all relevant images based on user query
        matching_results_image_fromdescription_data = utils.get_similar_image_from_query(
            st.session_state['text_df'],
            st.session_state['image_df'],
            query=prompt,
            image_query_path=st.session_state['image_uri'],
            column_name="embeddings",
            image_emb=False,
            top_n=10,
            embedding_size=1408,
            project_id=PROJECT_ID
        )

        context_text = []
        for key, value in matching_results_chunks_data.items():
            context_text.append(value["text"])
        final_context_text = "\n".join([str(item) for item in context_text])

        # combine all the relevant images and their description generated by Gemini
        context_images = []
        for key, value in matching_results_image_fromdescription_data.items():
            st.write(value["image_object"])
            st.write(value["image_description"])
            context_images.extend(
                ["Image: ", value["image_object"], "Caption: ", value["image_description"]]
            )

        final_prompt = f""" Instructions: The context of extraction of detail should be based on the text given in "final_context_text" : \n
        and image context in "context_images". Do not include any cumulative total return in the answer. ".

        Context:
        - Text Context:
        {final_context_text}
        - Image Context:
        {context_images}

        {prompt}

        Answer:
        """

        #instructions = """The context of extraction of detail should be based on the text context given in "text_context": \n
        #Base your response on "text_context". Do not include any cumulative total return in the answer. Context:
        #"""
        #final_prompt = [
        #    prompt,
        #    instructions,
        #    "text_context:",
        #    final_context_text
        #]

        response = multiturn_generate_content(
            model,
            final_prompt,
            generation_config=config,
        )

        with st.chat_message("assistant"):
            st.markdown(response)
            st.session_state.messages.append({"role": "assistant", "content": response})