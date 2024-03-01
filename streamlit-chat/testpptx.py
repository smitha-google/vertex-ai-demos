from pptx import Presentation

from pptx.util import Inches

prs = Presentation("./pages/Doc AI Prep.pptx")
for slide_number, slide in enumerate(prs.slides): 
    print(f"Slide {slide_number + 1}:") 
    for shape in slide.shapes: 
        if hasattr(shape, "text"):
            print(shape.text)